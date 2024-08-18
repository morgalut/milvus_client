"""
This module provides functions to extract text from various file formats such as DOCX, PDF, and PPTX.
"""

import os
from docx import Document
from pptx import Presentation
from pdfminer.high_level import extract_text

def extract_text_from_file(file_path):
    """
    Extracts text from a file based on its extension.
    
    Args:
        file_path (str): The path to the file from which to extract text.
    
    Returns:
        str: The extracted text, or None if the file type is unsupported.
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == ".docx":
        return extract_text_from_docx(file_path)

    if ext == ".pdf":
        return extract_text(file_path)

    if ext == ".pptx":
        return extract_text_from_pptx(file_path)

    return None

def extract_text_from_docx(file_path):
    """
    Extracts text from a DOCX file.
    
    Args:
        file_path (str): The path to the DOCX file.
    
    Returns:
        str: The extracted text.
    """
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    """
    Extracts text from a PPTX file.
    
    Args:
        file_path (str): The path to the PPTX file.
    
    Returns:
        str: The extracted text.
    """
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    return "\n".join(text_runs)
